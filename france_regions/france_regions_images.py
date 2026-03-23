"""
Script to download images about France regions
"""

import os
import re
import shutil
from argparse import ArgumentParser
from math import ceil

import imagesize
import requests
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from yaml import safe_load

CONFIG_FILE_DEFAULT = "france_regions_images.yaml"
IMAGES_DIR = "images_regions"

IMAGES_PER_SLIDE = 5
IMAGE_HEIGHT_WIDTH_RATIO = 0.75
SLIDE_HEIGHT_INCHES = 7.2
SLIDE_WIDTH_INCHES = 10
IMAGE_LEFT_OFFSET_DEFAULT = 0.6
IMAGE_TOP_OFFSET_DEFAULT = 1.7
LINE_INTERVAL_INCHES_DEFAULT = 0.8

IMAGE_MAX_PIXEL_HEIGHT = 600

SLD_LAYOUT_TITLE_SLIDE = 0
SLD_LAYOUT_TITLE_ONLY = 5

# Fonction pour télécharger une image depuis Unsplash (source libre)
def download_image(url, filename):
    """
    Download an image from a url

    :param url: image url
    :param filename: image filename
    :return: status (0 for success, http status code otherwhise)
    """
    response = requests.get(url)
    if response.status_code == 200:
        with open(filename, "wb") as f:
            f.write(response.content)
            status = 0
    else:
        print(f"Error downloading {url} to {filename}: {response.reason}")
        status = response.status_code

    return status


def main():
    parser = ArgumentParser()
    parser.add_argument(
        "--image-per-page",
        type=int,
        default=None,
        help="Number of images per page",
    )
    parser.add_argument(
        "--config",
        default=CONFIG_FILE_DEFAULT,
        help=f"Configuration file (D: {CONFIG_FILE_DEFAULT})",
    )
    parser.add_argument(
        "--no-download", action="store_true", default=False, help="Do not download images, use existing ones"
    )
    options = parser.parse_args()

    # Configuration file processing
    images_per_slide = IMAGES_PER_SLIDE
    images_per_region = None
    presentation_title = None
    image_bottom_alignment = False

    with open(options.config, "r", encoding="utf-8") as f:
        config = safe_load(f.read())

    if "regions_config" in config:
        with open(config["regions_config"], "r", encoding="utf-8") as f:
            regions_config = safe_load(f.read())

        if "regions" in regions_config:
            regions = regions_config["regions"]
        else:
            raise Exception(f"No regions found in {config["regions_config"]}")

        if "maps" in regions_config:
            maps = regions_config["maps"]
        else:
            maps = {}
    else:
        raise Exception(f"'regions_config' not found in configuration file {options.config}")

    if "layout" in config:
        if "region" in config["layout"]:
            if "max_images" in config["layout"]["region"]:
                images_per_region = config["layout"]["region"]["max_images"]
        if "slide" in config["layout"]:
            if "images" in config["layout"]["slide"]:
                images_per_slide = config["layout"]["slide"]["images"]
            if "image_alignment" in config["layout"]["slide"]:
                if config["layout"]["slide"]["image_alignment"] == "bottom":
                    image_bottom_alignment = True
                else:
                    raise Exception(
                        f"Invalid value found for layout/slide/image_alignment "
                        f"({config['layout']['slide']['image_alignment']})"
                    )
        if "title" in config["layout"]:
            presentation_title = config["layout"]["title"]

    # Command line options take precedence over config file
    if options.image_per_page:
        images_per_slide = options.image_per_page

    os.makedirs(IMAGES_DIR, exist_ok=True)
    image_paths = {}

    for region, places in regions.items():
        image_paths[region] = []
        for place, url in places.items():
            filename = f"{IMAGES_DIR}/{region}_{place}.jpg".replace(" ", "_")
            if options.no_download:
                if os.path.exists(filename):
                    print(f"{place}: using existing {filename}")
                else:
                    print(f"{filename} does not exist and --no-download option specified")
                    continue
            else:
                print(f"Téléchargement : {place}")
                if re.match(r"https*:", url):
                    status = download_image(url, filename)
                    if status != 0:
                        continue
                else:
                    # Assume it is a local file and copy it
                    shutil.copyfile(url, filename)
            image_paths[region].append({"place": place, "file": filename})

    # Create PowerPoint file and title slide
    prs = Presentation()
    # Create a title slide if a title is defined in configuration
    if presentation_title:
        slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_SLIDE]
        title_slide = prs.slides.add_slide(slide_layout)
        title = title_slide.shapes.title
        title.text = "Régions de France en image"

    if images_per_slide <= 3:
        lines = 1
    elif images_per_slide <= 6:
        lines = 2
    else:
        raise Exception("More than 6 images per slide currently not supported")
    images_per_line = ceil(images_per_slide / lines)

    # Pixels / inches ration
    max_image_height_inches = SLIDE_HEIGHT_INCHES - IMAGE_TOP_OFFSET_DEFAULT - LINE_INTERVAL_INCHES_DEFAULT
    inches_pixels_ratio = max_image_height_inches / IMAGE_MAX_PIXEL_HEIGHT

    # Compute image width based on number of images per line
    # full width includes left margin
    image_left_offset = IMAGE_LEFT_OFFSET_DEFAULT
    line_interval_inches = LINE_INTERVAL_INCHES_DEFAULT
    image_top_offset = IMAGE_TOP_OFFSET_DEFAULT
    image_full_width_inches = (SLIDE_WIDTH_INCHES - image_left_offset) / images_per_line
    image_width_inches = image_full_width_inches - image_left_offset
    image_height_inches = image_width_inches * IMAGE_HEIGHT_WIDTH_RATIO
    image_full_height_inches = image_height_inches + line_interval_inches

    total_height = image_full_height_inches * lines
    if total_height > SLIDE_HEIGHT_INCHES:
        # Adjust image width so that images fit in the slide height, based on height/width ratio
        # full height includes line interval
        image_full_height_inches = (SLIDE_HEIGHT_INCHES - image_top_offset) / lines
        image_height_inches = image_full_height_inches - line_interval_inches
        image_width_inches = image_height_inches / IMAGE_HEIGHT_WIDTH_RATIO
        image_full_width_inches = image_width_inches / IMAGE_HEIGHT_WIDTH_RATIO
        image_left_offset = (SLIDE_WIDTH_INCHES - (image_width_inches * images_per_line)) / (images_per_line + 1)
    else:
        # Center vertically the images
        max_height = SLIDE_HEIGHT_INCHES - image_top_offset
        line_interval_inches = (max_height - (image_height_inches * lines)) / (lines + 0.5)
        image_full_height_inches = image_height_inches + line_interval_inches
        image_top_offset += line_interval_inches * 0.5

    for region, images in image_paths.items():
        region_slide_num = 0

        # Compute left offset for last line to balance free space around images if the number of
        # images on this line is less than the number of images on other lines. It can be different
        # for each region, depending on the actual number of images for the region
        region_image_num = len(images)
        if images_per_region and len(images) > images_per_region:
            region_image_num = images_per_region
        region_lines = ceil(region_image_num / images_per_line)
        region_current_line = 0
        last_line_images = region_image_num % images_per_line
        if last_line_images == 0:
            last_line_left_offset = image_left_offset
        else:
            last_line_left_offset = (SLIDE_WIDTH_INCHES - (image_width_inches * last_line_images)) / (
                last_line_images + 1
            )

        if region in maps:
            width, height = imagesize.get(maps[region])
            max_image_height_inches = SLIDE_HEIGHT_INCHES - IMAGE_TOP_OFFSET_DEFAULT - LINE_INTERVAL_INCHES_DEFAULT
            height_inches = height * inches_pixels_ratio
            if height_inches > max_image_height_inches:
                slide_image_width = max_image_height_inches * width / height
            else:
                slide_image_width = SLIDE_WIDTH_INCHES - (2 * IMAGE_LEFT_OFFSET_DEFAULT)
            left_offset = (SLIDE_WIDTH_INCHES - slide_image_width) / 2
            slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_ONLY]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(
                maps[region],
                Inches(left_offset),
                Inches(IMAGE_TOP_OFFSET_DEFAULT),
                width=Inches(slide_image_width),
            )
            slide.shapes.title.text = region

        for i, image_params in enumerate(images):
            if images_per_region and i == images_per_region:
                break

            if (i % images_per_line) == 0:
                region_current_line += 1

            if (i % images_per_slide) == 0:
                region_slide_num += 1
                # Create a new slile
                slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_ONLY]
                slide = prs.slides.add_slide(slide_layout)
                # Titre
                title = slide.shapes.title
                slide_num_text = "" if region_slide_num == 1 else region_slide_num
                if slide_num_text:
                    slide_num_text = f" ({slide_num_text})"
                title.text = f"{region} {slide_num_text}"
                line = 0
            else:
                line = (i % images_per_slide) // images_per_line

            place = image_params["place"]
            image = image_params["file"]

            # Retrieve image size
            image_actual_width, image_actual_height = imagesize.get(image)
            image_scaling_factor = image_width_inches / (image_actual_width * inches_pixels_ratio)
            image_actual_height_inches = (image_actual_height * inches_pixels_ratio) * image_scaling_factor

            # Add image
            if region_current_line == region_lines:
                # Last line
                left_offset = last_line_left_offset
                if last_line_images == 0:
                    line_image_num = images_per_line
                else:
                    line_image_num = last_line_images
            else:
                left_offset = image_left_offset
                line_image_num = images_per_line
            left = left_offset + (i % line_image_num) * (image_width_inches + left_offset)
            top = image_top_offset + (image_full_height_inches * line)
            image_top = top
            if image_bottom_alignment and image_actual_height_inches <= image_height_inches:
                # Do not align on the bottom images whose are higher than the standard images
                image_top += image_height_inches - image_actual_height_inches
            slide.shapes.add_picture(image, Inches(left), Inches(image_top), width=Inches(image_width_inches))

            # Add caption, aligned vertically for all images, except when the image is too high
            if image_actual_height_inches < image_height_inches:
                image_actual_height_inches = image_height_inches
            caption = slide.shapes.add_textbox(
                Inches(left),
                Inches(top + image_actual_height_inches),
                Inches(image_width_inches),
                Inches(0.5),
            )
            p = caption.text_frame.paragraphs[0]
            p.text = place
            p.alignment = PP_ALIGN.CENTER

    # Sauvegarde
    prs.save(f"{IMAGES_DIR}/diaporama_regions_images.pptx")

    print("\n✅ Terminé ! Fichier créé : diaporama_regions_images.pptx")


if __name__ == "__main__":
    exit(main())
